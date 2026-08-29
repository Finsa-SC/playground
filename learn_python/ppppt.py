from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Define colors
DARK_BLUE = RGBColor(31, 59, 127)  # #1F3B7F
LIGHT_BLUE = RGBColor(74, 144, 226)  # #4A90E2
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(44, 62, 80)  # #2C3E50
LIGHT_GRAY = RGBColor(236, 240, 241)  # #ECF0F1
GREEN = RGBColor(39, 174, 96)  # #27AE60
RED = RGBColor(231, 76, 60)  # #E74C3C
TEXT_DARK = RGBColor(26, 26, 26)
TEXT_LIGHT = RGBColor(102, 102, 102)


def add_header_shape(slide, title_text):
    """Add header background and title"""
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.7)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_BLUE
    header_shape.line.color.rgb = DARK_BLUE

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.1),
        Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE


# ===== SLIDE 1: COVER =====
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_BLUE

# Main title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1))
tf = title_box.text_frame
tf.text = "PRESENTASI PKL"
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.8))
tf = subtitle_box.text_frame
tf.text = "Migrasi Peta GeoTab ke PostgreSQL/PostGIS"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = LIGHT_BLUE
p.alignment = PP_ALIGN.CENTER

# Line separator
line = slide1.shapes.add_connector(1, Inches(2), Inches(2.9), Inches(8), Inches(2.9))
line.line.color.rgb = LIGHT_BLUE
line.line.width = Pt(2)

# Student info
info_box = slide1.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(7), Inches(1.8))
tf = info_box.text_frame
tf.word_wrap = True
info_text = """Nama: Finsa Kusuma Putra
Kelas: 12 RPL 1
DUDI: PT Mahaka Digital Indonesia
Periode: 13 April - 7 Agustus 2026"""
tf.text = info_text
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = WHITE
    paragraph.alignment = PP_ALIGN.CENTER

# ===== SLIDE 2: ALAT & BAHAN =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_header_shape(slide2, "2. ALAT & BAHAN")

# Table data
tools_data = [
    ("Bahasa Pemrograman", "Python 3.x"),
    ("Library GIS", "GDAL/OGR2ogr, PostGIS"),
    ("Database", "PostgreSQL 12+"),
    ("Python Modules", "psycopg2-binary, python-dotenv, threading"),
    ("GIS Desktop", "QGIS"),
    ("Format File", "MapInfo TAB (.TAB, .DAT), SQL, CSV")
]

y_pos = 1.0
for idx, (category, items) in enumerate(tools_data):
    # Background color alternating
    row_shape = slide2.shapes.add_shape(
        1,
        Inches(0.5), Inches(y_pos),
        Inches(9), Inches(0.55)
    )
    row_shape.fill.solid()
    row_shape.fill.fore_color.rgb = LIGHT_GRAY if idx % 2 == 0 else WHITE
    row_shape.line.color.rgb = LIGHT_BLUE
    row_shape.line.width = Pt(1)

    # Category
    cat_box = slide2.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.05), Inches(3), Inches(0.45))
    tf = cat_box.text_frame
    tf.text = category
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE

    # Items
    items_box = slide2.shapes.add_textbox(Inches(4.0), Inches(y_pos + 0.05), Inches(5.3), Inches(0.45))
    tf = items_box.text_frame
    tf.word_wrap = True
    tf.text = items
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = TEXT_DARK

    y_pos += 0.65

# ===== SLIDE 3: LANGKAH KERJA =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_header_shape(slide3, "3. LANGKAH KERJA")

# Tugas 1 title
t1_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(4.5), Inches(0.35))
tf = t1_title.text_frame
tf.text = "TUGAS 1: Generate .sql dari .TAB"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = GREEN

# Tugas 1 steps
t1_steps = """1. Baca file .TAB dari folder
2. Parse geometri & validasi CRS
3. Transform ke EPSG:4326
4. Generate SQL INSERT
5. Export ke .sql file"""

t1_box = slide3.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(4), Inches(1.8))
tf = t1_box.text_frame
tf.text = t1_steps
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(12)
    paragraph.font.color.rgb = TEXT_DARK

# Tugas 2 title
t2_title = slide3.shapes.add_textbox(Inches(5.2), Inches(0.9), Inches(4.3), Inches(0.35))
tf = t2_title.text_frame
tf.text = "TUGAS 2: Insert ke PostgreSQL"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RED

# Tugas 2 steps
t2_steps = """1. Baca file .TAB dari folder
2. Parse & validasi geometri
3. Connect ke PostgreSQL
4. Insert ke 6 tabel (bidang,
   blok, bangunan, jalan,
   sungai, kelurahan)
5. Export database ke CSV/SQL"""

t2_box = slide3.shapes.add_textbox(Inches(5.4), Inches(1.3), Inches(4.1), Inches(1.8))
tf = t2_box.text_frame
tf.text = t2_steps
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(12)
    paragraph.font.color.rgb = TEXT_DARK

# Note
note_box = slide3.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.7))
tf = note_box.text_frame
tf.word_wrap = True
tf.text = "📌 Data: 6 tabel (map_bidang, map_blok, map_bangunan, map_jalan, map_sungai, map_kelurahan)"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = TEXT_LIGHT

# ===== SLIDE 4: KENDALA & SOLUSI =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_header_shape(slide4, "4. KENDALA & SOLUSI")

# Table headers
header_shape = slide4.shapes.add_shape(1, Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = LIGHT_BLUE
header_shape.line.color.rgb = LIGHT_BLUE

header1 = slide4.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(4), Inches(0.3))
tf = header1.text_frame
tf.text = "Kendala"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE

header2 = slide4.shapes.add_textbox(Inches(5.3), Inches(0.95), Inches(4), Inches(0.3))
tf = header2.text_frame
tf.text = "Solusi"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE

# Table data
kendala_data = [
    ("Library OGR2ogr asing & dokumentasi kompleks", "Belajar melalui contoh, trial & error, baca GDAL docs"),
    ("File .TAB tidak konsisten (CRS berbeda, format field beragam)",
     "Validasi CRS per file, transform ke EPSG:4326, field fallback"),
    ("Logika pengecekan NOP/geometri rumit", "Buat helper functions, validasi ST_IsValid(), skip error"),
    ("Baru pertama kali pakai QGIS & konsep GIS", "Inspect di QGIS, pahami proyeksi, geometri, atribut")
]

table_y = 1.4
for idx, (kendala, solusi) in enumerate(kendala_data):
    bg_color = LIGHT_GRAY if idx % 2 == 0 else WHITE

    # Kendala cell
    kendala_shape = slide4.shapes.add_shape(1, Inches(0.5), Inches(table_y), Inches(4.4), Inches(0.65))
    kendala_shape.fill.solid()
    kendala_shape.fill.fore_color.rgb = bg_color
    kendala_shape.line.color.rgb = LIGHT_BLUE
    kendala_shape.line.width = Pt(0.5)

    kendala_box = slide4.shapes.add_textbox(Inches(0.7), Inches(table_y + 0.05), Inches(4), Inches(0.55))
    tf = kendala_box.text_frame
    tf.word_wrap = True
    tf.text = kendala
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = TEXT_DARK

    # Solusi cell
    solusi_shape = slide4.shapes.add_shape(1, Inches(5.0), Inches(table_y), Inches(4.4), Inches(0.65))
    solusi_shape.fill.solid()
    solusi_shape.fill.fore_color.rgb = bg_color
    solusi_shape.line.color.rgb = LIGHT_BLUE
    solusi_shape.line.width = Pt(0.5)

    solusi_box = slide4.shapes.add_textbox(Inches(5.2), Inches(table_y + 0.05), Inches(4), Inches(0.55))
    tf = solusi_box.text_frame
    tf.word_wrap = True
    tf.text = solusi
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = TEXT_DARK

    table_y += 0.7

# ===== SLIDE 5: PEMBELAJARAN =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_header_shape(slide5, "5. PEMBELAJARAN & SKILL")

# Learning items
learnings = [
    ("📊", "SQL & Database Design", "Query kompleks, normalisasi, PostGIS functions, index"),
    ("🗺️", "QGIS & Konsep GIS", "Proyeksi (CRS), geometri, layer, atribut, validasi spatial"),
    ("🐍", "Python Advanced Logic", "Threading, GDAL/OGR, error handling, file parsing"),
    ("⚙️", "DevOps & Workflow", "Git commit, env variables, migration pipeline, QA")
]

learn_y = 1.0
for icon, title, desc in learnings:
    # Icon
    icon_box = slide5.shapes.add_textbox(Inches(0.6), Inches(learn_y), Inches(0.4), Inches(0.4))
    tf = icon_box.text_frame
    tf.text = icon
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide5.shapes.add_textbox(Inches(1.2), Inches(learn_y), Inches(3), Inches(0.25))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK_BLUE

    # Description
    desc_box = slide5.shapes.add_textbox(Inches(1.2), Inches(learn_y + 0.25), Inches(7.6), Inches(0.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    tf.text = desc
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = TEXT_LIGHT

    learn_y += 1.0

# Footer
footer_shape = slide5.shapes.add_shape(1, Inches(0.5), Inches(4.9), Inches(9), Inches(0.6))
footer_shape.fill.solid()
footer_shape.fill.fore_color.rgb = LIGHT_GRAY
footer_shape.line.color.rgb = LIGHT_GRAY

footer_box = slide5.shapes.add_textbox(Inches(0.7), Inches(4.95), Inches(8.6), Inches(0.5))
tf = footer_box.text_frame
tf.word_wrap = True
tf.text = "Terima kasih! Semoga project ini bermanfaat untuk pembelajaran GIS & data migration."
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = DARK_GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save
prs.save('/home/silence-suzuka/Presentasi_PKL_Finsa.pptx')
print("✅ PowerPoint berhasil dibuat: Presentasi_PKL_Finsa.pptx")